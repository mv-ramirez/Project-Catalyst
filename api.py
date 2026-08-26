import os
import io
import sys
import time
import json
import html
import base64
import asyncio
import logging
import threading
import requests
from contextlib import asynccontextmanager
from datetime import datetime
from urllib.parse import quote
from urllib3.util.retry import Retry
from requests.adapters import HTTPAdapter
from dotenv import load_dotenv
from docx import Document
import openpyxl
import pdfplumber
from pptx import Presentation
from fastapi import Depends, FastAPI, HTTPException, Query
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, HTMLResponse
from pydantic import BaseModel, Field
import uvicorn
import db

load_dotenv()

if sys.platform == "win32":
    asyncio.set_event_loop_policy(asyncio.WindowsSelectorEventLoopPolicy())

# Direct model cache to /tmp so it works on SAP BTP's Linux filesystem
os.environ.setdefault("SENTENCE_TRANSFORMERS_HOME", "/tmp/models")

from sentence_transformers import SentenceTransformer


# ── Logging ───────────────────────────────────────────────────────────────────

logging.basicConfig(
    level   = logging.INFO,
    format  = "%(asctime)s  %(levelname)-8s  %(name)s — %(message)s",
    datefmt = "%Y-%m-%d %H:%M:%S",
)
log = logging.getLogger("sharepoint_api")


# ── Config ────────────────────────────────────────────────────────────────────

_embed_model = None

def _get_embed_model():
    global _embed_model
    if _embed_model is None:
        _embed_model = SentenceTransformer("all-MiniLM-L6-v2")
    return _embed_model
_token_cache = {"token": None, "expires_at": 0}
_token_lock  = threading.Lock()

# Detect runtime environment
_ON_CF = bool(os.environ.get("VCAP_SERVICES"))

TENANT_ID     = os.getenv("TENANT_ID")
CLIENT_ID     = os.getenv("CLIENT_ID")
CLIENT_SECRET = os.getenv("CLIENT_SECRET")
HOSTNAME      = os.getenv("HOSTNAME",      "z0y8z.sharepoint.com")
SITE_PATH     = os.getenv("SITE_PATH",     "/sites/BTP-CatalystPlatform")
DOWNLOAD_PATH = os.getenv("DOWNLOAD_PATH", "/tmp")          # SAP BTP: use /tmp or Object Store mount
_log_day = datetime.now().strftime("%Y%m%d")
os.makedirs(os.path.join(DOWNLOAD_PATH, _log_day), exist_ok=True)
_fh = logging.FileHandler(os.path.join(DOWNLOAD_PATH, _log_day, f"api_{datetime.now().strftime('%Y%m%d_%H%M%S')}.log"), encoding="utf-8")
_fh.setFormatter(logging.Formatter(
    "%(asctime)s  %(levelname)-8s  %(name)s — %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
))
logging.getLogger().addHandler(_fh)


def _call_logger(endpoint: str) -> logging.Logger:
    ts      = datetime.now().strftime("%Y%m%d_%H%M%S")
    day     = ts[:8]
    log_dir = os.path.join(DOWNLOAD_PATH, day)
    os.makedirs(log_dir, exist_ok=True)
    path    = os.path.join(log_dir, f"{endpoint}_{ts}.log")
    logger  = logging.Logger(f"call.{endpoint}.{ts}")
    handler = logging.FileHandler(path, encoding="utf-8")
    handler.setFormatter(logging.Formatter(
        "%(asctime)s  %(levelname)-8s — %(message)s",
        datefmt="%Y-%m-%d %H:%M:%S",
    ))
    logger.addHandler(handler)
    logger.setLevel(logging.INFO)
    logger.propagate = False
    return logger

GRAPH_BASE = "https://graph.microsoft.com/v1.0"
SCOPE      = "https://graph.microsoft.com/.default"
TOKEN_URL  = f"https://login.microsoftonline.com/{TENANT_ID}/oauth2/v2.0/token"
URL_BASE   = f"{GRAPH_BASE}/sites/{HOSTNAME}:{SITE_PATH}"

site_id  = None
drive_id = None


# ── Auth ──────────────────────────────────────────────────────────────────────

# BTP Destination Service caches
_dest_svc_creds     = None
_dest_svc_token     = {"token": None, "expires_at": 0}
_dest_svc_lock      = threading.Lock()
DESTINATION_NAME    = "BTP-CP_MS_GRAPH"


def _get_dest_svc_creds() -> dict:
    """Extract the BTP Destination Service binding from VCAP_SERVICES."""
    global _dest_svc_creds
    if _dest_svc_creds:
        return _dest_svc_creds
    vcap = json.loads(os.environ.get("VCAP_SERVICES", "{}"))
    for key in ("destination", "destinations"):
        if key in vcap:
            _dest_svc_creds = vcap[key][0]["credentials"]
            log.info("BTP Destination Service binding found")
            return _dest_svc_creds
    raise RuntimeError(
        "Destination service not bound — add 'destination' under services: in manifest.yml"
    )


def _get_dest_svc_token() -> str:
    """Get an OAuth token scoped to the BTP Destination Service API."""
    with _dest_svc_lock:
        if _dest_svc_token["token"] and time.time() < _dest_svc_token["expires_at"] - 60:
            return _dest_svc_token["token"]
        creds = _get_dest_svc_creds()
        resp  = requests.post(
            f"{creds['url']}/oauth/token",
            data={
                "grant_type":    "client_credentials",
                "client_id":     creds["clientid"],
                "client_secret": creds["clientsecret"],
            },
            timeout=15,
        )
        resp.raise_for_status()
        data = resp.json()
        _dest_svc_token["token"]      = data["access_token"]
        _dest_svc_token["expires_at"] = time.time() + data.get("expires_in", 3600)
        log.info("BTP Destination Service token refreshed")
        return _dest_svc_token["token"]


def _get_ms_graph_token_from_destination() -> tuple[str, int]:
    """
    Fetch an MS Graph access token via the BTP Destination Service.
    Returns (token, expires_in_seconds).
    Prefers the pre-fetched authToken from the Destination Service;
    falls back to calling Azure AD directly using the destination's credentials.
    """
    creds      = _get_dest_svc_creds()
    dest_token = _get_dest_svc_token()
    resp = requests.get(
        f"{creds['uri']}/destination-configuration/v1/destinations/{DESTINATION_NAME}",
        headers={"Authorization": f"Bearer {dest_token}"},
        timeout=15,
    )
    resp.raise_for_status()
    payload = resp.json()

    # Option A — Destination Service pre-fetched a Bearer token
    for auth in payload.get("authTokens", []):
        if auth.get("type", "").lower() == "bearer" and auth.get("value"):
            expires_in = int(auth.get("expires_in", 3600))
            log.info(f"MS Graph token obtained from Destination Service authToken (expires_in={expires_in}s)")
            return auth["value"], expires_in

    # Option B — extract credentials and call Azure AD directly
    config     = payload.get("destinationConfiguration", {})
    token_url  = config.get("tokenServiceURL") or config.get("TokenServiceURL", "")
    client_id  = config.get("clientId")        or config.get("ClientId", "")
    client_sec = config.get("clientSecret")    or config.get("ClientSecret", "")
    scope      = config.get("scope",  SCOPE)
    if not (token_url and client_id and client_sec):
        raise RuntimeError(
            f"Destination '{DESTINATION_NAME}' missing tokenServiceURL / clientId / clientSecret"
        )
    resp2 = requests.post(
        token_url,
        data={
            "grant_type":    "client_credentials",
            "client_id":     client_id,
            "client_secret": client_sec,
            "scope":         scope,
        },
        timeout=15,
    )
    resp2.raise_for_status()
    data2 = resp2.json()
    expires_in = int(data2.get("expires_in", 3600))
    log.info(f"MS Graph token obtained via destination credentials (expires_in={expires_in}s)")
    return data2["access_token"], expires_in


def _getAccessToken() -> str:
    with _token_lock:
        if _token_cache["token"] and time.time() < _token_cache["expires_at"] - 60:
            return _token_cache["token"]

        if _ON_CF:
            token, expires_in = _get_ms_graph_token_from_destination()
        else:
            # Local dev — use env vars
            resp = requests.post(TOKEN_URL, data={
                "grant_type":    "client_credentials",
                "client_id":     CLIENT_ID,
                "client_secret": CLIENT_SECRET,
                "scope":         SCOPE,
            })
            resp.raise_for_status()
            data       = resp.json()
            token      = data["access_token"]
            expires_in = data.get("expires_in", 3600)

        _token_cache["token"]      = token
        _token_cache["expires_at"] = time.time() + expires_in
        return _token_cache["token"]

def _headers() -> dict:
    return {"Authorization": f"Bearer {_getAccessToken()}", "Accept": "application/json"}


# ── Site and Drive ────────────────────────────────────────────────────────────

def _get_SiteDrive_ID() -> tuple:
    r = requests.get(URL_BASE, headers=_headers(), timeout=15)
    r.raise_for_status()
    siteid = r.json()["id"]
    r.close()
    d = requests.get(f"{GRAPH_BASE}/sites/{siteid}/drives", headers=_headers(), timeout=15)
    d.raise_for_status()
    drives = d.json().get("value", [])
    if not drives:
        raise RuntimeError("No document libraries found in SharePoint site")
    driveid = drives[0]["id"]
    d.close()
    return siteid, driveid


# ── Formatters ────────────────────────────────────────────────────────────────

def _fmt_size(size_bytes: int) -> str:
    if size_bytes is None:
        return "—"
    for unit in ("B", "KB", "MB", "GB"):
        if size_bytes < 1024:
            return f"{size_bytes:.1f} {unit}"
        size_bytes /= 1024
    return f"{size_bytes:.1f} TB"


# ── Flatten helpers ───────────────────────────────────────────────────────────

def _flatten_item(item: dict) -> dict:
    ref      = item.get("parentReference", {})
    raw_path = ref.get("path", "")
    path     = raw_path.split(":")

    item_id     = item.get("id", "")
    item_name   = item.get("name", "")
    folder_path = path[1] if len(path) > 1 else ""
    dPath       = "/".join(quote(seg, safe="") for seg in f"{folder_path}/{item_name}".split("/"))
    basePath    = f"{GRAPH_BASE}{path[0]}:{dPath}:/content"

    return {
        "id":                 item_id,
        "name":               item_name,
        "type":               "Folder" if "folder" in item else item.get("file", {}).get("mimeType", "file"),
        "extension":          item.get("file", {}).get("fileExtension"),
        "size":               _fmt_size(item.get("size")),
        "size_bytes":         item.get("size"),
        "last_modified_date": item.get("lastModifiedDateTime", ""),
        "created_date":       item.get("createdDateTime", ""),
        "created_by":         item.get("createdBy", {}).get("user", {}).get("displayName", ""),
        "email_created_by":   item.get("createdBy", {}).get("user", {}).get("email", ""),
        "modified_by":        item.get("lastModifiedBy", {}).get("user", {}).get("displayName", ""),
        "email_modified_by":  item.get("lastModifiedBy", {}).get("user", {}).get("email", ""),
        "web_url":            item.get("webUrl", ""),
        "child_count":        item.get("folder", {}).get("childCount") if "folder" in item else None,
        "doc_path":           dPath,
        "versionURL":         f"{GRAPH_BASE}/drives/{drive_id}/items/{item_id}/versions",
        "BaseURL":            basePath,
        "shared":             item.get("shared", {}).get("scope", ""),
    }


def _flatten_versions(item: dict) -> dict:
    return {
        "version_id":           item.get("id"),
        "fileSize":             _fmt_size(item.get("size")),
        "size":                 item.get("size"),
        "lastModifiedDateTime": item.get("lastModifiedDateTime"),
        "lastmodified_by":      item.get("lastModifiedBy", {}).get("user", {}).get("email"),
        "download_url":         item.get("@microsoft.graph.downloadUrl"),
    }


# ── Session ───────────────────────────────────────────────────────────────────

def _build_session() -> requests.Session:
    session = requests.Session()
    retry   = Retry(
        total            = 3,
        backoff_factor   = 2,
        status_forcelist = [429, 500, 502, 503, 504],
    )
    session.mount("https://", HTTPAdapter(max_retries=retry))
    session.headers.update(_headers())
    return session


# ── Version fetcher ───────────────────────────────────────────────────────────

def _get_item_versions(id: str) -> list:
    rURL  = f"{GRAPH_BASE}/drives/{drive_id}/items/{id}/versions"
    rlist = []
    try:
        with _build_session() as session:
            res = session.get(rURL, timeout=(15, 120))
            res.raise_for_status()
            rlist = res.json().get("value", [])
            res.close()
            rlist = [_flatten_versions(x) for x in rlist]
            for x in rlist:
                x["item_id"]       = id
                x["version_count"] = len(rlist)
    except requests.exceptions.ReadTimeout:
        print(f"[TIMEOUT] item {id} — skipping")
    except requests.exceptions.HTTPError as e:
        sc = e.response.status_code if e.response is not None else "?"
        print(f"[HTTP ERROR] item {id} — {sc}")
    except Exception as e:
        print(f"[ERROR] item {id} — {e}")
    return rlist


# ── File listers ──────────────────────────────────────────────────────────────

def _get_root_files() -> list:
    log.info("Getting all root files from ~/root/delta")
    url    = f"{GRAPH_BASE}/drives/{drive_id}/root/delta"
    result = []
    first  = True
    while url:
        resp = requests.get(url, headers=_headers(), timeout=(15, 60))
        resp.raise_for_status()
        data = resp.json()
        resp.close()
        items = data.get("value", [])
        for f in (items[1:] if first else items):
            flat = _flatten_item(f)
            if flat["type"] != "Folder":
                result.append(flat)
        first = False
        url = data.get("@odata.nextLink")
    return result


def _get_all_files(src: str) -> list:
    log.info(f"Getting all files from /root{src}:/delta")
    url    = f"{GRAPH_BASE}/drives/{drive_id}/root{src}:/delta"
    result = []
    while url:
        resp = requests.get(url, headers=_headers(), timeout=(15, 60))
        resp.raise_for_status()
        data = resp.json()
        resp.close()
        for f in data.get("value", []):
            flat = _flatten_item(f)
            if flat["type"] != "Folder":
                result.append(flat)
        url = data.get("@odata.nextLink")
    return result


def _get_item_by_id(item_id: str) -> dict:
    url = f"{GRAPH_BASE}/drives/{drive_id}/items/{item_id}"
    r   = requests.get(url, headers=_headers(), timeout=15)
    r.raise_for_status()
    return _flatten_item(r.json())


def _search_files_by_name(query: str) -> list:
    url     = f"{GRAPH_BASE}/drives/{drive_id}/root/search(q='{query}')"
    results = []
    while url:
        r = requests.get(url, headers=_headers(), timeout=15)
        r.raise_for_status()
        data = r.json()
        for item in data.get("value", []):
            if "folder" not in item and query.lower() in item.get("name", "").lower():
                results.append(_flatten_item(item))
        url = data.get("@odata.nextLink")
    return results


def _get_root_files_with_versions(delay: float = 1.0) -> list:
    root_files = _get_root_files()
    total      = len(root_files)
    jdata      = []

    print(f"Found {total} files. Fetching versions...")
    log.info(f"Found {total} files. Fetching versions...")

    for i, x in enumerate(root_files, start=1):
        file_id  = x.get("id")
        doc_path = x.get("doc_path")
        print(f"  [{i}/{total}] {doc_path}")
        log.info(f"[{i}/{total}] {doc_path}")

        ver_list           = _get_item_versions(file_id)
        x["version_count"] = len(ver_list)
        x["versions_list"] = ver_list
        jdata.append(x)

        if i < total:
            time.sleep(delay)

    print(f"Done. {total} files processed.")
    log.info(f"Done. {total} files processed.")

    return jdata


# ── Version extractor ────────────────────────────────────────────────────────

def _split_versions(root_files: list) -> tuple:
    """
    Separates versions_list out of each file dict into a flat list.
    Annotates each version record with filename for traceability.
    version_count is kept in the file dict; versions_list is removed.
    Returns (clean_files, flat_versions).
    """
    clean_files   = []
    flat_versions = []
    for f in root_files:
        versions = f.pop("versions_list", None) or []
        filename = f.get("name", "")
        for v in versions:
            v["filename"] = filename
        flat_versions.extend(versions)
        clean_files.append(f)
    return clean_files, flat_versions


# ── JSON writers ──────────────────────────────────────────────────────────────

def _write_to_json(data: list, path: str = None) -> None:
    if path:
        with open(path, "a", encoding="utf-8") as fh:
            for record in data:
                fh.write(json.dumps(record, ensure_ascii=False) + "\n")
        print(f"\n{len(data)} records written to {path}")
        log.info(f"{len(data)} records written to {path}")


def _write_to_json_array_indent(data: list, path: str) -> None:
    with open(path, "w", encoding="utf-8") as fh:
        json.dump(data, fh, indent=2, ensure_ascii=False)
    print(f"\n{len(data)} records written to {path}")
    log.info(f"{len(data)} records written to {path}")


def _generate_analysis_html(run_summary: dict, save_dir: str) -> str:
    ts         = run_summary["ts"]
    total      = run_summary["total_items"]
    vectorized = run_summary["vectorized_count"]
    n_chunks   = run_summary["chunk_count"]
    skips      = run_summary["skips"]
    retries    = run_summary["retries"]
    file_sizes = run_summary["file_sizes"]
    duration_s = run_summary.get("duration_s", 0)

    no_text  = [s for s in skips if s["outcome"] == "no_text"]
    timeouts = [s for s in skips if s["outcome"] == "timeout_exhausted"]
    http_err = [s for s in skips if s["outcome"].startswith("http_")]
    other    = [s for s in skips if s["outcome"] in ("no_url", "error")]
    coverage = (vectorized / total * 100) if total else 0
    mins     = int(duration_s // 60)
    secs     = int(duration_s % 60)
    dur_str  = f"{mins}m {secs}s" if mins else f"{secs}s"

    def fmt_bytes(b):
        if not b: return "—"
        for u in ("B", "KB", "MB", "GB"):
            if b < 1024: return f"{b:.1f} {u}"
            b /= 1024
        return f"{b:.1f} TB"

    def skip_rows(items, color):
        return "".join(
            f'<tr><td>{s["filename"]}</td><td>{s["size"]}</td>'
            f'<td><span class="badge" style="background:{color}22;color:{color};border:1px solid {color}44">'
            f'{s["outcome"].replace("_"," ").upper()}</span></td></tr>'
            for s in items
        )

    retry_rows = "".join(
        f'<tr><td>{r["filename"]}</td><td>{r["size"]}</td>'
        f'<td style="color:{"#f97316" if r["attempts"]==3 else "#d8b4fe"};font-weight:700">{r["attempts"]}/3</td>'
        f'<td><span class="badge badge-green">recovered</span></td></tr>'
        for r in retries
    )

    skip_section = (
        (f'<h3 style="margin-top:1rem;margin-bottom:.5rem">Timeout Exhausted ({len(timeouts)})</h3>'
         f'<table><thead><tr><th>File</th><th>Size</th><th>Status</th></tr></thead>'
         f'<tbody>{skip_rows(timeouts,"#ef4444")}</tbody></table>' if timeouts else "") +
        (f'<h3 style="margin-top:1rem;margin-bottom:.5rem">HTTP Errors ({len(http_err)})</h3>'
         f'<table><thead><tr><th>File</th><th>Size</th><th>Status</th></tr></thead>'
         f'<tbody>{skip_rows(http_err,"#f97316")}</tbody></table>' if http_err else "") +
        (f'<h3 style="margin-top:1rem;margin-bottom:.5rem">No Text Extracted ({len(no_text)})</h3>'
         f'<table><thead><tr><th>File</th><th>Size</th><th>Status</th></tr></thead>'
         f'<tbody>{skip_rows(no_text,"#60a5fa")}</tbody></table>' if no_text else "") +
        (f'<h3 style="margin-top:1rem;margin-bottom:.5rem">Other ({len(other)})</h3>'
         f'<table><thead><tr><th>File</th><th>Size</th><th>Status</th></tr></thead>'
         f'<tbody>{skip_rows(other,"#94a3b8")}</tbody></table>' if other else "")
    )

    file_size_rows = "".join(
        f'<tr><td style="font-family:monospace;font-size:.78rem;color:#94a3b8">{n}</td>'
        f'<td style="font-weight:700;color:#e2e8f0">{fmt_bytes(s)}</td></tr>'
        for n, s in file_sizes.items()
    )

    retry_section = (
        f'<section><div class="section-title">Retry Detail ({len(retries)} files needed &gt;1 attempt)</div>'
        f'<div class="card card-orange"><table><thead><tr><th>File</th><th>Size</th>'
        f'<th>Attempts Used</th><th>Outcome</th></tr></thead>'
        f'<tbody>{retry_rows}</tbody></table></div></section>'
        if retries else ""
    )

    skip_color    = "card-orange" if (timeouts or http_err or other) else "card-green"
    retry_color   = "card-orange" if retries else "card-green"
    skip_val_col  = "orange" if (timeouts or http_err or other) else "green"
    retry_val_col = "orange" if retries else "green"

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<title>Vectorize Analysis — {ts}</title>
<style>
*,*::before,*::after{{box-sizing:border-box;margin:0;padding:0}}
body{{font-family:'Segoe UI',system-ui,sans-serif;background:#0f1117;color:#e2e8f0;padding:2rem;min-height:100vh}}
h1{{font-size:1.5rem;font-weight:700;color:#f8fafc;margin-bottom:.25rem}}
h2{{font-size:1rem;color:#94a3b8;font-weight:400;margin-bottom:.75rem}}
h3{{font-size:.85rem;font-weight:600;color:#cbd5e1;text-transform:uppercase;letter-spacing:.06em}}
.meta{{font-size:.75rem;color:#64748b;margin-top:.4rem;margin-bottom:1.75rem;padding-bottom:1rem;border-bottom:1px solid #1e293b}}
.grid-4{{display:grid;grid-template-columns:repeat(4,1fr);gap:1rem;margin-bottom:1.5rem}}
.card{{background:#1e293b;border:1px solid #334155;border-radius:10px;padding:1.25rem}}
.card-green{{border-color:#22c55e44}}.card-blue{{border-color:#3b82f644}}
.card-purple{{border-color:#a855f744}}.card-orange{{border-color:#f9731644}}
.kpi{{display:flex;flex-direction:column;gap:.2rem}}
.kpi .lbl{{font-size:.72rem;color:#94a3b8;text-transform:uppercase;letter-spacing:.06em}}
.kpi .val{{font-size:2rem;font-weight:700;line-height:1}}
.kpi .sub{{font-size:.78rem;margin-top:.2rem}}
.green{{color:#22c55e}}.blue{{color:#60a5fa}}.orange{{color:#fb923c}}.gray{{color:#94a3b8}}
.bar-wrap{{background:#0f1117;border-radius:6px;height:8px;overflow:hidden;margin-top:.4rem}}
.bar{{height:100%;border-radius:6px}}
.section-title{{font-size:.72rem;font-weight:700;text-transform:uppercase;letter-spacing:.08em;color:#475569;display:flex;align-items:center;gap:.5rem;margin-bottom:.75rem}}
.section-title::after{{content:'';flex:1;height:1px;background:#1e293b}}
section{{margin-bottom:1.75rem}}
table{{width:100%;border-collapse:collapse;font-size:.82rem}}
th{{padding:.55rem .8rem;text-align:left;color:#64748b;font-size:.72rem;font-weight:600;text-transform:uppercase;letter-spacing:.06em;border-bottom:1px solid #334155;background:#0f1117}}
td{{padding:.55rem .8rem;border-bottom:1px solid #1e293b;color:#cbd5e1;vertical-align:middle}}
tr:last-child td{{border-bottom:none}}
tr:hover td{{background:#1e293b55}}
.badge{{display:inline-block;padding:.12rem .5rem;border-radius:999px;font-size:.7rem;font-weight:600}}
.badge-green{{background:#14532d55;color:#4ade80;border:1px solid #16a34a55}}
</style>
</head>
<body>
<h1>Vectorize Analysis</h1>
<h2>SharePoint RAG API · Auto-generated run report</h2>
<div class="meta">
  Run timestamp: <code>{ts}</code> &nbsp;|&nbsp; Duration: <strong>{dur_str}</strong> &nbsp;|&nbsp; Files scanned: {total}
</div>

<section>
<div class="section-title">Summary</div>
<div class="grid-4">
  <div class="card card-green">
    <div class="kpi">
      <span class="lbl">Files Vectorized</span>
      <span class="val green">{vectorized}</span>
      <span class="sub gray">{coverage:.1f}% of {total}</span>
      <div class="bar-wrap"><div class="bar" style="width:{coverage:.1f}%;background:linear-gradient(90deg,#16a34a,#22c55e)"></div></div>
    </div>
  </div>
  <div class="card card-blue">
    <div class="kpi">
      <span class="lbl">Chunk Records</span>
      <span class="val blue">{n_chunks:,}</span>
      <span class="sub gray">semantic chunks embedded</span>
    </div>
  </div>
  <div class="card {skip_color}">
    <div class="kpi">
      <span class="lbl">Total Skipped</span>
      <span class="val {skip_val_col}">{len(skips)}</span>
      <span class="sub gray">{len(no_text)} no-text · {len(timeouts)} timeout · {len(http_err)} HTTP err</span>
    </div>
  </div>
  <div class="card {retry_color}">
    <div class="kpi">
      <span class="lbl">Retries Triggered</span>
      <span class="val {retry_val_col}">{len(retries)}</span>
      <span class="sub gray">{"files >1 attempt — all recovered" if retries else "all succeeded on attempt 1"}</span>
    </div>
  </div>
</div>
</section>

{retry_section}

<section>
<div class="section-title">Skipped Files ({len(skips)} total)</div>
<div class="card">
{skip_section or '<p style="color:#64748b;font-size:.82rem;padding:.5rem 0">No files skipped — all files downloaded and extracted successfully.</p>'}
</div>
</section>

<section>
<div class="section-title">Output Files</div>
<div class="card" style="max-width:640px">
<table><thead><tr><th>File</th><th>Size</th></tr></thead>
<tbody>{file_size_rows}</tbody></table>
</div>
</section>

<div class="meta" style="margin-top:2rem;text-align:center">
  Auto-generated by pytesting/catalyst/for_SAP/api.py &nbsp;·&nbsp; {ts}
</div>
</body>
</html>"""

    path = os.path.join(save_dir, f"sp_analysis_{ts}.html")
    with open(path, "w", encoding="utf-8") as fh:
        fh.write(html)
    log.info(f"Analysis HTML written to {path}")
    return path


# ── Chunk JSONL reader ────────────────────────────────────────────────────────

def _read_chunks_from_jsonl(path: str, filename: str) -> list:
    log.info(f"Reading chunks from {path}")
    chunks = []
    try:
        with open(path, encoding="utf-8") as fh:
            for line in fh:
                line = line.strip()
                if not line:
                    continue
                try:
                    record = json.loads(line)
                    if record.get("filename") == filename:
                        chunks.append(record)
                except json.JSONDecodeError:
                    continue
    except FileNotFoundError:
        pass

    log.info(f"Returning {len(chunks)} chunks for {filename}")
    return chunks


# ── File downloader (bytes + base64 chunks) ───────────────────────────────────

def _download_file_content(item: dict, save_dir: str = None) -> bytes:
    url       = item.get("BaseURL")
    file_name = item.get("name")

    if not url:
        log.warning(f"No BaseURL for {file_name}")
        return b""

    if save_dir:
        os.makedirs(save_dir, exist_ok=True)
        json_path = os.path.join(save_dir, "files_list_chunk.json")
        if not os.path.exists(json_path):
            with open(json_path, "w"):
                pass
        log.info(f"Saving chunks to {json_path}")

    print(f"Downloading: {file_name} ({item.get('size', '—')})")
    log.info(f"Downloading: {file_name} ({item.get('size', '—')})")

    try:
        chunks = []
        with requests.get(url, headers=_headers(), stream=True, timeout=(15, 300)) as r:
            r.raise_for_status()
            for i, chunk in enumerate(r.iter_content(chunk_size=8192), start=1):
                if chunk:
                    chunks.append(chunk)
                    encoded    = base64.b64encode(chunk).decode("utf-8")
                    chunk_list = [{
                        "item_id":        item.get("id"),
                        "filename":       item.get("name"),
                        "chunk_datetime": datetime.now().isoformat(),
                        "chunk_size":     len(chunk),
                        "chunk_id":       i,
                        "chunk_data":     encoded,
                    }]
                    if save_dir:
                        _write_to_json(chunk_list, path=os.path.join(save_dir, "files_list_chunk.json"))
                    print(f"  chunk [{i}]: {len(chunk)} bytes | {encoded[:20]}...")

        total_chunks = len(chunks)
        content      = b"".join(chunks)
        print(f"  Downloaded {len(content) / 1024:.1f} KB  ({total_chunks} chunks)")
        log.info(f"  Downloaded {len(content) / 1024:.1f} KB  ({total_chunks} chunks)")

        if save_dir:
            save_path = os.path.join(save_dir, file_name)
            with open(save_path, "wb") as f:
                f.write(content)
            print(f"  Saved to: {save_path}")
            log.info(f"  Saved to: {save_path}")

        return content

    except requests.exceptions.ReadTimeout:
        print(f"[TIMEOUT] {file_name} — skipping")
        log.warning(f"[TIMEOUT] {file_name} — skipping")
    except requests.exceptions.HTTPError as e:
        sc = e.response.status_code if e.response is not None else "?"
        print(f"[HTTP ERROR] {file_name} — {sc}")
        log.error(f"[HTTP ERROR] {file_name} — {sc}")
    except Exception as e:
        print(f"[ERROR] {file_name} — {e}")
        log.error(f"[ERROR] {file_name} — {e}")
    return b""


# ── Chunk rebuilder ───────────────────────────────────────────────────────────

def _rebuild_file(chunks: list, save_dir: str = None) -> bytes:
    if not chunks:
        print("[ERROR] No chunks provided")
        log.error("No chunks provided")
        return b""

    try:
        sorted_chunks = sorted(chunks, key=lambda x: x.get("chunk_id", 0))
        file_name     = sorted_chunks[0].get("filename")
        total         = len(sorted_chunks)

        print(f"Rebuilding: {file_name} ({total} chunks)")
        log.info(f"Rebuilding: {file_name} ({total} chunks)")

        decoded = []
        for c in sorted_chunks:
            chunk_id = c.get("chunk_id")
            original = base64.b64decode(c.get("chunk_data", "").encode("utf-8"))
            decoded.append(original)
            print(f"  chunk [{chunk_id}]: {len(original)} bytes decoded")
            log.info(f"  chunk [{chunk_id}]: {len(original)} bytes decoded")

        content = b"".join(decoded)
        print(f"  Rebuilt {len(content) / 1024:.1f} KB from {total} chunks")
        log.info(f"  Rebuilt {len(content) / 1024:.1f} KB from {total} chunks")

        if save_dir:
            os.makedirs(save_dir, exist_ok=True)
            save_path = os.path.join(save_dir, file_name)
            with open(save_path, "wb") as f:
                f.write(content)
            print(f"  Saved to: {save_path}")
            log.info(f"  Saved to: {save_path}")

        return content

    except Exception as e:
        print(f"[ERROR] Rebuild failed — {e}")
        log.error(f"[ERROR] Rebuild failed — {e}")
        return b""


# ── RAG: text extraction, chunking, vectorization ────────────────────────────

def _extract_text(content: bytes, filename: str) -> str:
    log.info(f"Extracting text from {filename}")
    ext = filename.rsplit(".", 1)[-1].lower()
    try:
        if ext in ("docx", "dotx"):
            log.info("Extracting docx/dotx.")
            doc = Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        elif ext == "doc":
            log.info("Extracting doc (binary Word).")
            import olefile, re
            with olefile.OleFileIO(io.BytesIO(content)) as ole:
                if ole.exists("WordDocument"):
                    raw  = ole.openstream("WordDocument").read()
                    text = raw.decode("utf-16-le", errors="ignore")
                    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", text)
                    text = re.sub(r" {3,}", " ", text).strip()
                    return text
            return ""

        elif ext in ("xlsx", "xlsm"):
            log.info("Extracting xlsx.")
            try:
                wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True)
            except Exception as e:
                if "zip" in str(e).lower():
                    log.warning("  xlsx is not a zip file — retrying with xlrd (likely legacy XLS binary)")
                    import xlrd
                    wb = xlrd.open_workbook(file_contents=content)
                    lines = []
                    for sheet in wb.sheets():
                        for row_idx in range(sheet.nrows):
                            line = " | ".join(
                                str(sheet.cell_value(row_idx, col))
                                for col in range(sheet.ncols)
                                if sheet.cell_value(row_idx, col) not in (None, "")
                            )
                            if line.strip():
                                lines.append(line)
                    return "\n".join(lines)
                raise
            lines = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    line = " | ".join(str(c) for c in row if c is not None)
                    if line.strip():
                        lines.append(line)
            return "\n".join(lines)

        elif ext == "xls":
            log.info("Extracting xls (legacy binary).")
            import xlrd
            wb    = xlrd.open_workbook(file_contents=content)
            lines = []
            for sheet in wb.sheets():
                for row_idx in range(sheet.nrows):
                    line = " | ".join(
                        str(sheet.cell_value(row_idx, col))
                        for col in range(sheet.ncols)
                        if sheet.cell_value(row_idx, col) not in (None, "")
                    )
                    if line.strip():
                        lines.append(line)
            return "\n".join(lines)

        elif ext == "pdf":
            log.info("Extracting pdf.")
            text = []
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    text.append(page.extract_text() or "")
            return "\n".join(text)

        elif ext in ("pptx", "pptm"):
            log.info("Extracting pptx/pptm.")
            prs   = Presentation(io.BytesIO(content))
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                lines.append(t)
            return "\n".join(lines)

        elif ext == "ppt":
            log.info("Extracting ppt (binary PowerPoint).")
            import olefile, re
            texts = []
            with olefile.OleFileIO(io.BytesIO(content)) as ole:
                for entry in ole.listdir():
                    try:
                        raw     = ole.openstream(entry).read()
                        decoded = raw.decode("utf-16-le", errors="ignore")
                        cleaned = re.sub(r"[\x00-\x1f]", " ", decoded).strip()
                        if len(cleaned) > 20:
                            texts.append(cleaned)
                    except Exception:
                        continue
            return "\n".join(texts)

    except Exception as e:
        print(f"[ERROR] Text extraction failed for {filename} — {e}")
        log.error(f"[ERROR] Text extraction failed for {filename} — {e}")
    return ""


def _semantic_chunks(text: str, max_tokens: int = 512) -> list:
    max_chars  = max_tokens * 4
    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
    chunks     = []
    current    = []
    length     = 0

    log.info(f"Chunking into ~{max_tokens} token segments")

    for para in paragraphs:
        if length + len(para) > max_chars and current:
            chunks.append("\n".join(current))
            current = []
            length  = 0
        current.append(para)
        length += len(para)

    if current:
        chunks.append("\n".join(current))

    log.info(f"Chunked into {len(chunks)} segments")

    return chunks


def _download_and_vectorize(items: list[dict], logger: logging.Logger = None) -> tuple:
    _log            = logger or log
    ts              = datetime.now().strftime("%Y%m%d_%H%M%S")
    chunks_data     = []
    raw_chunks_data = []
    vectors_data    = []

    t_start = time.time()
    skips   = []
    retries = []

    _log.info(f"Starting vectorization of {len(items)} items")

    for idx, item in enumerate(items, start=1):
        if idx > 1:
            time.sleep(1.0)

        url      = item.get("BaseURL")
        filename = item.get("name", "")
        item_id  = item.get("id")

        print(f"[{idx}/{len(items)}] Downloading: {filename} ({item.get('size', '—')})")
        _log.info(f"Downloading: {filename} ({item.get('size', '—')})")

        if not url:
            print(f"  [ERROR] No BaseURL — skipping")
            _log.error(f"  [ERROR] No BaseURL — skipping")
            skips.append({"filename": filename, "size": item.get("size", "—"), "attempts": 0, "outcome": "no_url"})
            continue

        _attempts    = 0
        _outcome     = None
        MAX_ATTEMPTS = 3
        BACKOFF      = [5, 15, 30]
        content      = None
        raw_chunks   = []

        for attempt in range(1, MAX_ATTEMPTS + 1):
            _attempts = attempt
            try:
                _log.info(f"Downloading {filename} (attempt {attempt}/{MAX_ATTEMPTS}) via {url.split('?')[0]}...")
                raw_chunks  = []
                byte_chunks = []
                with requests.get(url, headers=_headers(), stream=True, timeout=(15, 300)) as r:
                    r.raise_for_status()
                    for i, chunk in enumerate(r.iter_content(chunk_size=8192), start=1):
                        if chunk:
                            byte_chunks.append(chunk)
                            raw_chunks.append({
                                "chunk_id":   i,
                                "chunk_size": len(chunk),
                                "chunk_data": base64.b64encode(chunk).decode("utf-8"),
                            })
                content = b"".join(byte_chunks)
                print(f"  Downloaded {len(content) / 1024:.1f} KB ({len(raw_chunks)} raw chunks)")
                _log.info(f"  Downloaded {len(content) / 1024:.1f} KB ({len(raw_chunks)} raw chunks)")
                _outcome = "ok"
                break

            except requests.exceptions.ReadTimeout:
                print(f"  [TIMEOUT] attempt {attempt}/{MAX_ATTEMPTS}")
                _log.warning(f"  [TIMEOUT] attempt {attempt}/{MAX_ATTEMPTS}")
                if attempt < MAX_ATTEMPTS:
                    time.sleep(BACKOFF[attempt - 1])
                else:
                    print(f"  [TIMEOUT] all attempts exhausted — skipping")
                    _log.warning(f"  [TIMEOUT] all attempts exhausted — skipping")
                    _outcome = "timeout_exhausted"

            except requests.exceptions.HTTPError as e:
                sc = e.response.status_code if e.response is not None else 0
                if sc == 429 and attempt < MAX_ATTEMPTS:
                    retry_after = int(e.response.headers.get("Retry-After", BACKOFF[attempt - 1]))
                    print(f"  [429] Retry-After {retry_after}s — waiting...")
                    _log.warning(f"  [429] Retry-After {retry_after}s — waiting")
                    time.sleep(retry_after)
                elif sc in (500, 502, 503, 504) and attempt < MAX_ATTEMPTS:
                    print(f"  [HTTP {sc}] attempt {attempt}/{MAX_ATTEMPTS} — retrying in {BACKOFF[attempt - 1]}s...")
                    _log.warning(f"  [HTTP {sc}] attempt {attempt}/{MAX_ATTEMPTS} — retrying in {BACKOFF[attempt - 1]}s")
                    time.sleep(BACKOFF[attempt - 1])
                else:
                    print(f"  [HTTP ERROR] {sc} — {e.response.text[:120]}")
                    _log.error(f"  [HTTP ERROR] {sc}")
                    _outcome = f"http_{sc}"
                    break

            except Exception as e:
                print(f"  [ERROR] {e}")
                _log.error(f"  [ERROR] {e}")
                _outcome = "error"
                break

        if content is None:
            skips.append({"filename": filename, "size": item.get("size", "—"), "attempts": _attempts, "outcome": _outcome or "error"})
            continue

        if _attempts > 1:
            retries.append({"filename": filename, "size": item.get("size", "—"), "attempts": _attempts, "outcome": "recovered"})

        text = _extract_text(content, filename)
        if not text.strip():
            print(f"  [SKIP] No text extracted")
            _log.warning(f"  [SKIP] No text extracted")
            skips.append({"filename": filename, "size": item.get("size", "—"), "attempts": _attempts, "outcome": "no_text"})
            continue

        chunks = _semantic_chunks(text, max_tokens=512)
        print(f"  Chunking → {len(chunks)} semantic chunks")
        _log.info(f"  Chunking → {len(chunks)} semantic chunks")

        for i, chunk_text in enumerate(chunks, start=1):
            embedding = _get_embed_model().encode(chunk_text).tolist()

            chunks_data.append({
                "item_id":   item_id,
                "filename":  filename,
                "chunk_id":  i,
                "chunk_data": chunk_text,
            })

            vectors_data.append({
                "item_id":    item_id,
                "filename":   filename,
                "chunk_id":   i,
                "vector_id":  i,
                "vector_data": embedding,
            })

        raw_chunks_data.append({
            "item_id":     item_id,
            "filename":    filename,
            "chunk_count": len(chunks),
            "raw_chunks":  raw_chunks,
        })

        print(f"  {len(chunks)} chunks vectorized ({len(chunks[-1]) if chunks else 0} chars last)")
        _log.info(f"  {len(chunks)} chunks vectorized ({len(chunks[-1]) if chunks else 0} chars last)")

    print(f"\n[BATCH COMPLETE] {len(chunks_data)} chunks, {len(raw_chunks_data)} files vectorized.")
    _log.info(f"Batch complete — {len(chunks_data)} chunk records, {len(raw_chunks_data)} file records.")

    active_dir   = ts[:8]
    save_dir     = os.path.join(DOWNLOAD_PATH, active_dir)
    os.makedirs(save_dir, exist_ok=True)
    path_chunks  = os.path.join(save_dir, f"sp_chunks_{ts}.json")
    path_raw     = os.path.join(save_dir, f"sp_files_chunk_{ts}.json")
    path_vectors = os.path.join(save_dir, f"sp_vectors_{ts}.json")
    _write_to_json_array_indent(chunks_data,     path=path_chunks)
    _write_to_json_array_indent(raw_chunks_data, path=path_raw)
    _write_to_json_array_indent(vectors_data,    path=path_vectors)

    path_html = _generate_analysis_html(
        run_summary={
            "ts":               ts,
            "total_items":      len(items),
            "vectorized_count": len(raw_chunks_data),
            "chunk_count":      len(chunks_data),
            "skips":            skips,
            "retries":          retries,
            "duration_s":       time.time() - t_start,
            "file_sizes": {
                f"sp_chunks_{ts}.json":      os.path.getsize(path_chunks),
                f"sp_files_chunk_{ts}.json": os.path.getsize(path_raw),
                f"sp_vectors_{ts}.json":     os.path.getsize(path_vectors),
            },
        },
        save_dir=save_dir,
    )
    return chunks_data, raw_chunks_data, vectors_data, path_chunks, path_raw, path_vectors, path_html


# ── Lifespan ──────────────────────────────────────────────────────────────────

@asynccontextmanager
async def lifespan(app: FastAPI):
    global site_id, drive_id

    # Validate credentials source before attempting any network call
    if _ON_CF:
        try:
            _get_dest_svc_creds()
            log.info(f"CF mode — MS Graph credentials via BTP Destination '{DESTINATION_NAME}'")
        except RuntimeError as e:
            log.error(str(e))
            raise
    else:
        missing = [v for v in ("TENANT_ID", "CLIENT_ID", "CLIENT_SECRET") if not os.getenv(v)]
        if missing:
            log.error(f"Missing required env vars: {', '.join(missing)}")
            raise RuntimeError(f"Missing required env vars: {', '.join(missing)}")

    log.info("Connecting to SharePoint...")
    try:
        site_id, drive_id = await asyncio.to_thread(_get_SiteDrive_ID)
        log.info(f"Connected — site_id={site_id}  drive_id={drive_id}")
    except Exception as e:
        log.warning(f"SharePoint init failed (non-fatal): {e} — endpoints will return 503 until connection is available")

    try:
        table_status = await asyncio.to_thread(db.create_tables)
        for table, result in table_status.items():
            log.info(f"  DB table {table}: {result}")
    except Exception as e:
        log.warning(f"HANA table setup failed (non-fatal): {e}")

    yield
    log.info("Shutdown complete.")


# ── App ───────────────────────────────────────────────────────────────────────

app = FastAPI(
    title       = "SharePoint RAG API",
    version     = "1.0.0",
    description = "Fetch, download, vectorize, and rebuild SharePoint documents.",
    lifespan    = lifespan,
)

app.add_middleware(
    CORSMiddleware,
    allow_origins  = ["*"],
    allow_methods  = ["*"],
    allow_headers  = ["*"],
)

app.include_router(db.router)


# ── Dependency ────────────────────────────────────────────────────────────────

def require_connection():
    if drive_id is None:
        raise HTTPException(status_code=503, detail="SharePoint connection not ready")


# ── Models ────────────────────────────────────────────────────────────────────

class VectorizeRequest(BaseModel):
    items:    list[dict] | None = Field(None,  description="Items to vectorize. Omit to auto-fetch root files with versions.")
    save_dir: str        | None = Field(None,  description="Output directory. Defaults to DOWNLOAD_PATH/<YYYYMMDD>.")

class VectorizeResponse(BaseModel):
    count_chunks: int
    count_files:  int
    saved_to:     list[str]

class PipelineRequest(BaseModel):
    save_dir: str | None = Field(None, description="Output directory. Defaults to DOWNLOAD_PATH/<YYYYMMDD>.")
    delay:    float      = Field(1.0,  ge=0.0, le=10.0, description="Seconds between version requests.")

class PipelineResponse(BaseModel):
    count_root_files: int
    count_chunks:     int
    count_files:      int
    saved_to:         list[str]

class DownloadRequest(BaseModel):
    item:     dict
    save_dir: str | None = None

class RebuildRequest(BaseModel):
    chunks:   list[dict]
    save_dir: str | None = None

class VectorizeItemRequest(BaseModel):
    item_id: str = Field(..., description="SharePoint item ID to download and vectorize")

class VectorizeItemResponse(BaseModel):
    item_id:      str
    filename:     str
    count_chunks: int
    saved_to:     list[str]


# ── Routes ────────────────────────────────────────────────────────────────────

@app.get("/health", tags=["system"])
def health():
    clog = _call_logger("health")
    clog.info("GET /health called")
    result = {"status": "ok", "site_id": site_id, "drive_id": drive_id}
    clog.info(f"Response: {result}")
    return result


@app.get("/debug/destination", tags=["system"])
def debug_destination():
    """
    Diagnostic endpoint — checks if the BTP Destination Service can read 'btp-cp_ms_graph'.
    Returns the raw destination payload (minus clientSecret) or a detailed error.
    Only available when running on CF (_ON_CF=True).
    """
    if not _ON_CF:
        return {"on_cf": False, "message": "Not running on CF — destination service not applicable"}

    result: dict = {"on_cf": True, "destination_name": DESTINATION_NAME}

    # Step 1: get Destination Service binding
    try:
        creds = _get_dest_svc_creds()
        result["dest_svc_uri"]      = creds.get("uri")
        result["dest_svc_url"]      = creds.get("url")
        result["dest_svc_clientid"] = creds.get("clientid")
        result["step_binding"]      = "ok"
    except Exception as e:
        result["step_binding"] = f"FAILED: {e}"
        return result

    # Step 2: get XSUAA token for Destination Service
    try:
        dest_token = _get_dest_svc_token()
        result["step_xsuaa_token"] = "ok"
    except Exception as e:
        result["step_xsuaa_token"] = f"FAILED: {e}"
        return result

    # Step 3: call Destination Service API for the named destination
    try:
        resp = requests.get(
            f"{creds['uri']}/destination-configuration/v1/destinations/{DESTINATION_NAME}",
            headers={"Authorization": f"Bearer {dest_token}"},
            timeout=15,
        )
        result["step_dest_api_status"] = resp.status_code
        if resp.status_code == 200:
            payload = resp.json()
            config  = payload.get("destinationConfiguration", {})
            # Strip secret before returning
            config_safe = {k: v for k, v in config.items() if "secret" not in k.lower() and "Secret" not in k}
            result["step_dest_api"]       = "ok"
            result["destination_config"]  = config_safe
            result["auth_tokens_present"] = [
                {"type": t.get("type"), "has_value": bool(t.get("value"))}
                for t in payload.get("authTokens", [])
            ]
        else:
            result["step_dest_api"] = f"FAILED: HTTP {resp.status_code}"
            result["response_body"] = resp.text[:500]
    except Exception as e:
        result["step_dest_api"] = f"FAILED: {e}"

    # Step 4: try getting an MS Graph token end-to-end
    try:
        token, expires_in = _get_ms_graph_token_from_destination()
        result["step_ms_graph_token"] = f"ok — expires_in={expires_in}s, token_prefix={token[:20]}..."
    except Exception as e:
        result["step_ms_graph_token"] = f"FAILED: {e}"

    return result


@app.get("/debug/db", tags=["system"])
def debug_db():
    """
    Diagnostic endpoint — checks HANA database connectivity and reports
    the host, port, schema, and table status.
    """
    result: dict = {}

    # Step 1: connect
    try:
        conn = db._connect()
        result["step_connect"] = "ok"
        result["schema"]       = db._schema or "(none)"
        result["connected"]    = conn.isconnected()
    except Exception as e:
        result["step_connect"] = f"FAILED: {e}"
        return result

    # Step 2: check which tables exist
    try:
        cur = conn.cursor()
        cur.execute("SELECT TABLE_NAME FROM SYS.TABLES WHERE SCHEMA_NAME = CURRENT_SCHEMA")
        existing = {row[0] for row in cur.fetchall()}
        cur.close()

        expected = [
            "CATALYSTPLATFORM_BRAIN_FILES",
            "CATALYSTPLATFORM_BRAIN_TEXT_CHUNKS",
            "CATALYSTPLATFORM_BRAIN_VECTORS",
            "CATALYSTPLATFORM_BRAIN_FILES_VERSIONS",
            "CATALYSTPLATFORM_BRAIN_FILES_CHUNK",
        ]
        result["tables"] = {
            t: ("exists" if t in existing else "missing")
            for t in expected
        }
        result["step_tables"] = "ok"
    except Exception as e:
        result["step_tables"] = f"FAILED: {e}"

    return result


@app.get("/analysis", tags=["analysis"])
def list_analysis():
    """Lists all analysis HTML reports and log files in dated output folders, newest first."""
    results = []
    try:
        day_dirs = sorted(os.listdir(DOWNLOAD_PATH), reverse=True)
    except FileNotFoundError:
        return {"count": 0, "reports": []}
    for day_dir in day_dirs:
        day_path = os.path.join(DOWNLOAD_PATH, day_dir)
        if not os.path.isdir(day_path):
            continue
        for fname in sorted(os.listdir(day_path), reverse=True):
            if fname.startswith("sp_analysis_") and fname.endswith(".html"):
                results.append({"date": day_dir, "type": "analysis", "filename": fname,
                                 "path": os.path.join(day_path, fname), "url": f"/analysis/{day_dir}/{fname}"})
            elif fname.endswith(".log"):
                results.append({"date": day_dir, "type": "log", "filename": fname,
                                 "path": os.path.join(day_path, fname), "url": f"/analysis/{day_dir}/{fname}"})
    return {"count": len(results), "reports": results}


@app.get("/analysis/latest", tags=["analysis"], response_class=HTMLResponse)
def analysis_latest():
    """Returns the most recent vectorize analysis HTML report rendered in the browser."""
    try:
        day_dirs = sorted(os.listdir(DOWNLOAD_PATH), reverse=True)
    except FileNotFoundError:
        raise HTTPException(status_code=404, detail="No analysis reports found. Run /vectorize first.")
    for day_dir in day_dirs:
        day_path = os.path.join(DOWNLOAD_PATH, day_dir)
        if not os.path.isdir(day_path):
            continue
        for fname in sorted(os.listdir(day_path), reverse=True):
            if fname.startswith("sp_analysis_") and fname.endswith(".html"):
                with open(os.path.join(day_path, fname), encoding="utf-8") as fh:
                    return HTMLResponse(content=fh.read())
    raise HTTPException(status_code=404, detail="No analysis reports found. Run /vectorize first.")


@app.get("/analysis/{day}/{filename}", tags=["analysis"], response_class=HTMLResponse)
def analysis_by_name(day: str, filename: str):
    """Returns a specific analysis HTML or log file from a dated output folder."""
    safe_root = os.path.realpath(DOWNLOAD_PATH)
    resolved  = os.path.realpath(os.path.join(DOWNLOAD_PATH, day, filename))
    if not resolved.startswith(safe_root + os.sep):
        raise HTTPException(status_code=400, detail="Invalid path")
    if not os.path.isfile(resolved):
        raise HTTPException(status_code=404, detail="File not found")
    if filename.endswith(".html"):
        with open(resolved, encoding="utf-8") as fh:
            return HTMLResponse(content=fh.read())
    if filename.endswith(".log"):
        with open(resolved, encoding="utf-8") as fh:
            raw = fh.read()
        safe_name = html.escape(filename)
        safe_body = html.escape(raw)
        return HTMLResponse(content=(
            f"<!DOCTYPE html><html><head><meta charset='UTF-8'/>"
            f"<title>{safe_name}</title></head>"
            f"<body style='background:#0f1117;color:#e2e8f0;font-family:monospace;padding:2rem'>"
            f"<h2 style='color:#94a3b8;margin-bottom:1rem'>{safe_name}</h2>"
            f"<pre style='font-size:.82rem;line-height:1.6;white-space:pre-wrap'>{safe_body}</pre>"
            f"</body></html>"
        ))
    raise HTTPException(status_code=400, detail="Only .html and .log files are supported")


@app.get("/files", dependencies=[Depends(require_connection)], tags=["files"])
def list_files():
    clog = _call_logger("files")
    clog.info("GET /files called")
    try:
        result = _get_root_files()
        clog.info(f"Returned {len(result)} files")
        return result
    except requests.HTTPError as e:
        clog.error(f"HTTP error: {e.response.status_code}")
        raise HTTPException(status_code=e.response.status_code, detail=str(e))
    except Exception as e:
        clog.error(f"Error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/files/versions", dependencies=[Depends(require_connection)], tags=["files"])
def list_files_with_versions(
    delay: float = Query(1.0, ge=0.0, le=10.0, description="Seconds between version requests"),
):
    clog = _call_logger("files_versions")
    clog.info(f"GET /files/versions called — delay={delay}")
    try:
        result = _get_root_files_with_versions(delay=delay)
        clog.info(f"Returned {len(result)} files with versions")
        return result
    except requests.HTTPError as e:
        clog.error(f"HTTP error: {e.response.status_code}")
        raise HTTPException(status_code=e.response.status_code, detail=str(e))
    except Exception as e:
        clog.error(f"Error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/files/all", dependencies=[Depends(require_connection)], tags=["files"])
def list_all_files(
    path: str = Query("", description="Relative path from drive root, e.g. /Documents/Subfolder. Omit for root."),
):
    clog = _call_logger("files_all")
    clog.info(f"GET /files/all called — path='{path}'")
    try:
        result = _get_root_files() if not path else _get_all_files(path)
        clog.info(f"Returned {len(result)} files")
        return result
    except requests.HTTPError as e:
        clog.error(f"HTTP error: {e.response.status_code}")
        raise HTTPException(status_code=e.response.status_code, detail=str(e))
    except Exception as e:
        clog.error(f"Error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/files/search", dependencies=[Depends(require_connection)], tags=["files"])
def search_files(
    name: str = Query(..., description="Filename or partial name to search for"),
):
    clog = _call_logger("files_search")
    clog.info(f"GET /files/search called — name='{name}'")
    try:
        results = _search_files_by_name(name)
        clog.info(f"Found {len(results)} matches for '{name}'")
        return {"query": name, "count": len(results), "items": results}
    except requests.HTTPError as e:
        status = e.response.status_code if e.response is not None else 500
        clog.error(f"HTTP error: {status}")
        raise HTTPException(status_code=status, detail=str(e))
    except Exception as e:
        clog.error(f"Search error: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get(
    "/vectorize",
    response_model = PipelineResponse,
    dependencies   = [Depends(require_connection)],
    tags           = ["rag"],
    summary        = "Fetch all root files, download, and vectorize in one call",
)
def get_vectorize(
    save_dir: str   | None = Query(None, description="Output directory. Defaults to DOWNLOAD_PATH/<YYYYMMDD>."),
    delay:    float        = Query(1.0,  ge=0.0, le=10.0, description="Seconds between version requests."),
):
    clog = _call_logger("vectorize")
    clog.info(f"GET /vectorize called — delay={delay}, save_dir={save_dir}")
    log.info("GET /vectorize — fetching root files with versions")
    try:
        root_files = _get_root_files_with_versions(delay=delay)
    except Exception as e:
        clog.error(f"Failed to fetch root files: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch root files: {e}")

    if not root_files:
        clog.warning("No files found in SharePoint root")
        raise HTTPException(status_code=404, detail="No files found in SharePoint root")

    root_files, versions_data = _split_versions(root_files)
    clog.info(f"Fetched {len(root_files)} root files ({len(versions_data)} version records) — starting vectorization")
    try:
        chunks_data, raw_chunks_data, vectors_data, path_chunks, path_raw, path_vectors, path_html = _download_and_vectorize(root_files, logger=clog)
    except Exception as e:
        clog.error(f"Vectorization failed: {e}")
        raise HTTPException(status_code=500, detail=f"Vectorization failed: {e}")

    ts           = datetime.now().strftime("%Y%m%d_%H%M%S")
    save_dir     = save_dir or os.path.join(DOWNLOAD_PATH, ts[:8])
    os.makedirs(save_dir, exist_ok=True)
    path_root     = os.path.join(save_dir, f"sp_root_data_{ts}.json")
    path_versions = os.path.join(save_dir, f"sp_files_versions_{ts}.json")
    _write_to_json_array_indent(root_files,    path=path_root)
    _write_to_json_array_indent(versions_data, path=path_versions)

    log.info(f"GET /vectorize complete — {len(root_files)} files, {len(chunks_data)} chunks → {save_dir}")
    clog.info(f"Complete — {len(root_files)} files, {len(chunks_data)} chunks saved to {save_dir}")

    try:
        db.ingest(chunks_data, vectors_data, root_files)
        db.ingest_chunk_summary(raw_chunks_data)
        db.ingest_versions(versions_data)
        clog.info("DB ingest complete")
    except Exception as e:
        clog.warning(f"DB ingest failed (non-fatal): {e}")

    return PipelineResponse(
        count_root_files = len(root_files),
        count_chunks     = len(chunks_data),
        count_files      = len(raw_chunks_data),
        saved_to         = [path_chunks, path_raw, path_vectors, path_root, path_versions, path_html],
    )


@app.post(
    "/vectorize",
    response_model = VectorizeResponse,
    dependencies   = [Depends(require_connection)],
    tags           = ["rag"],
    summary        = "Download and vectorize a supplied list of items (POST with body)",
)
def vectorize(req: VectorizeRequest):
    clog = _call_logger("vectorize_post")
    clog.info(f"POST /vectorize called — items={len(req.items or [])}, save_dir={req.save_dir}")
    items         = req.items or []
    versions_data = []

    if not items:
        log.info("No items supplied — fetching root files with versions")
        clog.info("No items supplied — fetching root files with versions")
        try:
            fetched = _get_root_files_with_versions()
            items, versions_data = _split_versions(fetched)
        except Exception as e:
            clog.error(f"Failed to fetch root files: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to fetch root files: {e}")

    if not items:
        clog.warning("No files found to vectorize")
        raise HTTPException(status_code=404, detail="No files found to vectorize")

    clog.info(f"Vectorizing {len(items)} items")
    try:
        chunks_data, raw_chunks_data, vectors_data, path_chunks, path_raw, path_vectors, path_html = _download_and_vectorize(items, logger=clog)
    except Exception as e:
        clog.error(f"Vectorization failed: {e}")
        raise HTTPException(status_code=500, detail=f"Vectorization failed: {e}")

    saved_to = [path_chunks, path_raw, path_vectors, path_html]
    if versions_data:
        path_versions = os.path.join(os.path.dirname(path_chunks), f"sp_files_versions_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json")
        _write_to_json_array_indent(versions_data, path=path_versions)
        saved_to.append(path_versions)

    log.info(f"Vectorize complete — {len(chunks_data)} chunks, {len(raw_chunks_data)} files → {path_chunks}")
    clog.info(f"Complete — {len(chunks_data)} chunks, {len(raw_chunks_data)} files saved")

    try:
        db.ingest(chunks_data, vectors_data, items)
        db.ingest_chunk_summary(raw_chunks_data)
        db.ingest_versions(versions_data)
        clog.info("DB ingest complete")
    except Exception as e:
        clog.warning(f"DB ingest failed (non-fatal): {e}")

    return VectorizeResponse(
        count_chunks = len(chunks_data),
        count_files  = len(raw_chunks_data),
        saved_to     = saved_to,
    )


@app.post(
    "/vectorize/item",
    response_model = VectorizeItemResponse,
    dependencies   = [Depends(require_connection)],
    tags           = ["rag"],
    summary        = "Fetch a single SharePoint item by item_id and vectorize it",
)
def vectorize_single_item(req: VectorizeItemRequest):
    clog = _call_logger("vectorize_item")
    clog.info(f"POST /vectorize/item called — item_id={req.item_id}")
    try:
        item = _get_item_by_id(req.item_id)
    except requests.HTTPError as e:
        status = e.response.status_code if e.response is not None else 500
        clog.error(f"Item lookup failed: {status}")
        raise HTTPException(status_code=status, detail=f"Item not found: {req.item_id}")
    except Exception as e:
        clog.error(f"Item lookup error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

    clog.info(f"Fetched item: {item['name']}")
    try:
        chunks_data, raw_chunks_data, vectors_data, path_chunks, path_raw, path_vectors, path_html = _download_and_vectorize([item], logger=clog)
    except Exception as e:
        clog.error(f"Vectorization failed: {e}")
        raise HTTPException(status_code=500, detail=f"Vectorization failed: {e}")

    try:
        db.ingest(chunks_data, vectors_data, [item])
        db.ingest_chunk_summary(raw_chunks_data)
        clog.info("DB ingest complete")
    except Exception as e:
        clog.warning(f"DB ingest failed (non-fatal): {e}")

    return VectorizeItemResponse(
        item_id      = req.item_id,
        filename     = item["name"],
        count_chunks = len(chunks_data),
        saved_to     = [path_chunks, path_raw, path_vectors, path_html],
    )


@app.post(
    "/pipeline",
    response_model = PipelineResponse,
    dependencies   = [Depends(require_connection)],
    tags           = ["rag"],
    summary        = "Full pipeline: fetch root files → download → vectorize → save all three JSONs",
)
def pipeline(req: PipelineRequest = PipelineRequest()):
    clog = _call_logger("pipeline")
    clog.info(f"POST /pipeline called — delay={req.delay}, save_dir={req.save_dir}")
    log.info("Pipeline started — fetching root files with versions")
    try:
        root_files, versions_data = _split_versions(_get_root_files_with_versions(delay=req.delay))
    except Exception as e:
        clog.error(f"Failed to fetch root files: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch root files: {e}")

    if not root_files:
        clog.warning("No files found in SharePoint root")
        raise HTTPException(status_code=404, detail="No files found in SharePoint root")

    clog.info(f"Fetched {len(root_files)} root files ({len(versions_data)} version records) — starting vectorization")
    try:
        chunks_data, raw_chunks_data, vectors_data, path_chunks, path_raw, path_vectors, path_html = _download_and_vectorize(root_files, logger=clog)
    except Exception as e:
        clog.error(f"Vectorization failed: {e}")
        raise HTTPException(status_code=500, detail=f"Vectorization failed: {e}")

    ts            = datetime.now().strftime("%Y%m%d_%H%M%S")
    save_dir      = req.save_dir or os.path.join(DOWNLOAD_PATH, ts[:8])
    os.makedirs(save_dir, exist_ok=True)
    path_root     = os.path.join(save_dir, f"sp_root_data_{ts}.json")
    path_versions = os.path.join(save_dir, f"sp_files_versions_{ts}.json")
    _write_to_json_array_indent(root_files,    path=path_root)
    _write_to_json_array_indent(versions_data, path=path_versions)

    log.info(f"Pipeline complete — {len(root_files)} files, {len(chunks_data)} chunks → {save_dir}")
    clog.info(f"Complete — {len(root_files)} files, {len(chunks_data)} chunks saved to {save_dir}")

    try:
        db.ingest(chunks_data, vectors_data, root_files)
        db.ingest_chunk_summary(raw_chunks_data)
        db.ingest_versions(versions_data)
        clog.info("DB ingest complete")
    except Exception as e:
        clog.warning(f"DB ingest failed (non-fatal): {e}")

    return PipelineResponse(
        count_root_files = len(root_files),
        count_chunks     = len(chunks_data),
        count_files      = len(raw_chunks_data),
        saved_to         = [path_chunks, path_raw, path_vectors, path_root, path_versions, path_html],
    )


@app.post("/download", dependencies=[Depends(require_connection)], tags=["files"])
def download_file(req: DownloadRequest):
    filename = req.item.get("name", "file")
    clog = _call_logger("download")
    clog.info(f"POST /download called — file={filename}, save_dir={req.save_dir}")
    content = _download_file_content(req.item, save_dir=req.save_dir)
    if not content:
        clog.error(f"File not found or download failed: {filename}")
        raise HTTPException(status_code=404, detail="File not found or download failed")
    clog.info(f"Downloaded {len(content) / 1024:.1f} KB — streaming back to caller")
    return StreamingResponse(
        io.BytesIO(content),
        media_type = "application/octet-stream",
        headers    = {"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@app.get("/rebuild", tags=["files"])
def rebuild_file_get(
    filename: str       = Query(...,  description="Filename to rebuild, e.g. report.pdf"),
    save_dir: str | None = Query(None, description="Directory containing files_list_chunk.json. Defaults to DOWNLOAD_PATH/<today>."),
):
    """
    Reconstructs a previously downloaded file from the chunk records saved by
    POST /download. Reads files_list_chunk.json in save_dir, filters to the
    requested filename, reassembles the bytes, and streams the file back.
    Use this when you want to re-serve a file without hitting SharePoint again.
    """
    clog = _call_logger("rebuild")
    clog.info(f"GET /rebuild called — filename={filename}, save_dir={save_dir}")
    active_dir  = datetime.now().strftime("%Y%m%d")
    resolved    = save_dir or os.path.join(DOWNLOAD_PATH, active_dir)
    chunk_path  = os.path.join(resolved, "files_list_chunk.json")

    chunks = _read_chunks_from_jsonl(chunk_path, filename)
    if not chunks:
        clog.error(f"No chunks found for '{filename}' in {chunk_path}")
        raise HTTPException(
            status_code = 404,
            detail      = f"No chunks found for '{filename}' in {chunk_path}",
        )

    clog.info(f"Found {len(chunks)} chunks — rebuilding")
    content = _rebuild_file(chunks, save_dir=save_dir)
    if not content:
        clog.error("Rebuild failed — check chunk data")
        raise HTTPException(status_code=400, detail="Rebuild failed — check chunk data")

    clog.info(f"Rebuilt {len(content) / 1024:.1f} KB — streaming back to caller")
    return StreamingResponse(
        io.BytesIO(content),
        media_type = "application/octet-stream",
        headers    = {"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@app.post("/rebuild", tags=["files"])
def rebuild_file(req: RebuildRequest):
    clog = _call_logger("rebuild_post")
    clog.info(f"POST /rebuild called — chunks={len(req.chunks)}, save_dir={req.save_dir}")
    if not req.chunks:
        clog.error("chunks list is empty")
        raise HTTPException(status_code=422, detail="chunks list is empty")
    content = _rebuild_file(req.chunks, save_dir=req.save_dir)
    if not content:
        clog.error("Rebuild failed — check chunk data")
        raise HTTPException(status_code=400, detail="Rebuild failed — check chunk data")
    filename = req.chunks[0].get("filename", "file")
    clog.info(f"Rebuilt {len(content) / 1024:.1f} KB for {filename} — streaming back to caller")
    return StreamingResponse(
        io.BytesIO(content),
        media_type = "application/octet-stream",
        headers    = {"Content-Disposition": f'attachment; filename="{filename}"'},
    )

#final
# ── Entry point ───────────────────────────────────────────────────────────────

if __name__ == "__main__":
    uvicorn.run("api:app", host="0.0.0.0", port=int(os.getenv("PORT", 8000)), reload=False, timeout_keep_alive=600)
